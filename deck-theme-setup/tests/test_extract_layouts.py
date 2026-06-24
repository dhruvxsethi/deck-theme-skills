import os
import json
import shutil
import tempfile
import unittest

from pptx import Presentation
from scripts.extract_layouts import (
    classify_slide_role,
    is_valid_pptx,
    list_template_layouts,
    run_setup,
    learn_role_mapping,
    run_learn,
)


class TestClassifySlideRole(unittest.TestCase):
    def test_matches_keywords(self):
        self.assertEqual(classify_slide_role("Agenda", 1, 8), "agenda")
        self.assertEqual(classify_slide_role("Thank You", 7, 8), "thank-you")
        self.assertEqual(classify_slide_role("Next Steps", 6, 8), "next-steps")
        self.assertEqual(classify_slide_role("Current Challenges", 2, 8), "problem")
        self.assertEqual(classify_slide_role("Our Solution", 3, 8), "solution")
        self.assertEqual(classify_slide_role("Customer Success Story", 4, 8), "case-study")
        self.assertEqual(classify_slide_role("Return on Investment", 5, 8), "roi")

    def test_position_fallback(self):
        self.assertEqual(classify_slide_role("Welcome", 0, 8), "title")
        self.assertEqual(classify_slide_role("Goodbye", 7, 8), "thank-you")

    def test_content_fallback(self):
        self.assertEqual(classify_slide_role("Platform Overview", 3, 8), "content")


class TestTemplateInspection(unittest.TestCase):
    def _default_template(self):
        fd, path = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        Presentation().save(path)
        return path

    def test_is_valid_pptx_true_for_real_pptx(self):
        path = self._default_template()
        self.addCleanup(os.remove, path)
        self.assertTrue(is_valid_pptx(path))

    def test_is_valid_pptx_false_for_non_zip(self):
        fd, path = tempfile.mkstemp(suffix=".pptx")
        os.write(fd, b"\xe5\xe4\x88\xdb not a zip")
        os.close(fd)
        self.addCleanup(os.remove, path)
        self.assertFalse(is_valid_pptx(path))

    def test_list_template_layouts_returns_standard_names(self):
        path = self._default_template()
        self.addCleanup(os.remove, path)
        names = list_template_layouts(path)
        self.assertIn("Title Slide", names)
        self.assertIn("Title and Content", names)


class TestRunSetup(unittest.TestCase):
    def setUp(self):
        self.dir = tempfile.mkdtemp()
        self.addCleanup(shutil.rmtree, self.dir)
        self.template = os.path.join(self.dir, "src-template.pptx")
        Presentation().save(self.template)
        self.config = os.path.join(self.dir, "deck-config.json")
        self.templates_dir = os.path.join(self.dir, "templates")

    def test_setup_writes_variant_and_copies_template(self):
        variant = run_setup(self.template, "brand-light", "Brand — Light",
                            self.config, self.templates_dir)
        self.assertEqual(variant["id"], "brand-light")
        self.assertIn("Title Slide", variant["layouts"])
        self.assertEqual(variant["rolePreferences"], {})
        copied = os.path.join(self.templates_dir, "brand-light.pptx")
        self.assertTrue(os.path.exists(copied))
        with open(self.config) as f:
            cfg = json.load(f)
        self.assertEqual(len(cfg["variants"]), 1)
        self.assertEqual(cfg["variants"][0]["id"], "brand-light")

    def test_setup_replaces_same_id_variant(self):
        run_setup(self.template, "brand-light", "First", self.config, self.templates_dir)
        run_setup(self.template, "brand-light", "Second", self.config, self.templates_dir)
        with open(self.config) as f:
            cfg = json.load(f)
        self.assertEqual(len(cfg["variants"]), 1)
        self.assertEqual(cfg["variants"][0]["name"], "Second")

    def test_setup_rejects_non_pptx(self):
        bad = os.path.join(self.dir, "bad.pptx")
        with open(bad, "wb") as f:
            f.write(b"\xe5\xe4\x88\xdb")
        with self.assertRaisesRegex(ValueError, "doesn't look like a real"):
            run_setup(bad, "x", "X", self.config, self.templates_dir)


class TestLearn(unittest.TestCase):
    def _example_deck(self, dir_):
        prs = Presentation()
        layouts = {l.name: l for l in prs.slide_layouts}
        s1 = prs.slides.add_slide(layouts["Title Slide"])
        s1.shapes.title.text = "Welcome to Acme"
        s2 = prs.slides.add_slide(layouts["Title and Content"])
        s2.shapes.title.text = "Current Challenges"
        path = os.path.join(dir_, "example.pptx")
        prs.save(path)
        return path

    def test_learn_maps_role_to_actual_layout(self):
        d = tempfile.mkdtemp()
        self.addCleanup(shutil.rmtree, d)
        mapping = learn_role_mapping(self._example_deck(d))
        self.assertEqual(mapping["title"], "Title Slide")
        self.assertEqual(mapping["problem"], "Title and Content")

    def test_run_learn_merges_into_existing_variant(self):
        d = tempfile.mkdtemp()
        self.addCleanup(shutil.rmtree, d)
        template = os.path.join(d, "t.pptx")
        Presentation().save(template)
        config = os.path.join(d, "deck-config.json")
        run_setup(template, "v1", "V1", config, os.path.join(d, "templates"))
        run_learn(self._example_deck(d), "v1", config)
        with open(config) as f:
            cfg = json.load(f)
        prefs = cfg["variants"][0]["rolePreferences"]
        self.assertEqual(prefs["title"], "Title Slide")
        self.assertEqual(prefs["problem"], "Title and Content")

    def test_run_learn_errors_when_variant_missing(self):
        d = tempfile.mkdtemp()
        self.addCleanup(shutil.rmtree, d)
        config = os.path.join(d, "deck-config.json")
        with open(config, "w") as f:
            json.dump({"variants": []}, f)
        with self.assertRaisesRegex(ValueError, 'No variant "v1"'):
            run_learn(self._example_deck(d), "v1", config)


if __name__ == "__main__":
    unittest.main()
