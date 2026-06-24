"""deck-theme-setup: list a template's layouts and learn role->layout mappings."""

import json
import os
import re
import shutil

from pptx import Presentation

ROLE_KEYWORDS = [
    ("agenda", [r"agenda", r"what we.?ll cover"]),
    ("next-steps", [r"next steps"]),
    ("thank-you", [r"thank you", r"questions\?"]),
    ("problem", [r"challenge", r"problem", r"pain point", r"current state"]),
    ("solution", [r"solution", r"how we help", r"our approach"]),
    ("case-study", [r"case study", r"success stor", r"customer story"]),
    ("roi", [r"\broi\b", r"return on investment", r"business case"]),
]


def classify_slide_role(title, slide_index, total_slides):
    text = title or ""
    for role, patterns in ROLE_KEYWORDS:
        if any(re.search(p, text, re.IGNORECASE) for p in patterns):
            return role
    if slide_index == 0:
        return "title"
    if slide_index == total_slides - 1:
        return "thank-you"
    return "content"


def is_valid_pptx(path):
    """A .pptx/.potx is a zip; its first two bytes are the 'PK' signature."""
    with open(path, "rb") as f:
        return f.read(2) == b"PK"


def list_template_layouts(template_path):
    """Return the layout names from the template's first slide master."""
    prs = Presentation(template_path)
    return [layout.name for layout in prs.slide_layouts]


def _load_config(config_path):
    if os.path.exists(config_path):
        with open(config_path) as f:
            config = json.load(f)
    else:
        config = {"variants": []}
    if not isinstance(config.get("variants"), list):
        raise ValueError(f'Config at "{config_path}" is missing a valid "variants" array.')
    return config


def run_setup(template_path, variant_id, variant_name, config_path, templates_dir):
    if not is_valid_pptx(template_path):
        raise ValueError(
            f'"{template_path}" doesn\'t look like a real .pptx/.potx file (no zip '
            f"signature). It may have been encrypted or corrupted in transit — re-obtain it."
        )
    os.makedirs(templates_dir, exist_ok=True)
    dest = os.path.join(templates_dir, f"{variant_id}.pptx")
    shutil.copyfile(template_path, dest)

    config = _load_config(config_path)
    config["variants"] = [v for v in config["variants"] if v.get("id") != variant_id]
    config_dir = os.path.dirname(os.path.abspath(config_path))
    variant = {
        "id": variant_id,
        "name": variant_name,
        "template": os.path.relpath(dest, config_dir),
        "layouts": list_template_layouts(dest),
        "rolePreferences": {},
    }
    config["variants"].append(variant)
    with open(config_path, "w") as f:
        json.dump(config, f, indent=2)
    return variant


def learn_role_mapping(example_deck_path):
    prs = Presentation(example_deck_path)
    slides = list(prs.slides)
    total = len(slides)
    mapping = {}
    for i, slide in enumerate(slides):
        title = ""
        if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text or ""
        role = classify_slide_role(title, i, total)
        # Last slide with a given role wins — later slides override earlier ones.
        mapping[role] = slide.slide_layout.name
    return mapping


def run_learn(example_deck_path, variant_id, config_path):
    if not os.path.exists(config_path):
        raise ValueError(f'No config at "{config_path}" — run "setup" first.')
    config = _load_config(config_path)
    variant = next((v for v in config["variants"] if v.get("id") == variant_id), None)
    if variant is None:
        raise ValueError(f'No variant "{variant_id}" in {config_path} — run "setup" first.')
    mapping = learn_role_mapping(example_deck_path)
    variant["rolePreferences"] = {**variant.get("rolePreferences", {}), **mapping}
    with open(config_path, "w") as f:
        json.dump(config, f, indent=2)
    return mapping


def _main(argv=None):
    import argparse

    parser = argparse.ArgumentParser(prog="extract_layouts")
    sub = parser.add_subparsers(dest="mode", required=True)

    s = sub.add_parser("setup")
    s.add_argument("template")
    s.add_argument("variant_id")
    s.add_argument("variant_name")
    s.add_argument("config")
    s.add_argument("--templates-dir", default=None)

    l = sub.add_parser("learn")
    l.add_argument("example")
    l.add_argument("variant_id")
    l.add_argument("config")

    args = parser.parse_args(argv)
    if args.mode == "setup":
        templates_dir = args.templates_dir or os.path.join(
            os.path.dirname(os.path.abspath(args.config)), "templates"
        )
        variant = run_setup(args.template, args.variant_id, args.variant_name,
                            args.config, templates_dir)
        print(f'Saved variant "{variant["name"]}" with {len(variant["layouts"])} layouts:')
        for name in variant["layouts"]:
            print(f"  - {name}")
    elif args.mode == "learn":
        mapping = run_learn(args.example, args.variant_id, args.config)
        print(f'Learned role mappings for variant "{args.variant_id}":')
        for role, layout in mapping.items():
            print(f"  {role} -> {layout}")


if __name__ == "__main__":
    import sys

    try:
        _main()
    except Exception as exc:  # noqa: BLE001 - surface a clean message to the CLI user
        print(str(exc), file=sys.stderr)
        sys.exit(1)
