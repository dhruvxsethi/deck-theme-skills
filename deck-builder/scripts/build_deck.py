"""deck-builder: render a deck by populating a real template's own layouts."""

import json
import os

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE

BODY_TYPES = (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT)


def load_config(config_path):
    with open(config_path) as f:
        config = json.load(f)
    if not isinstance(config.get("variants"), list) or not config["variants"]:
        raise ValueError(f'Config at "{config_path}" has no variants.')
    return config


def get_variant(config, variant_id):
    for variant in config["variants"]:
        if variant.get("id") == variant_id:
            return variant
    ids = ", ".join(v.get("id", "?") for v in config["variants"])
    raise ValueError(f'No variant "{variant_id}" in config. Available: {ids}')


def _layout_has_body(layout):
    return any(ph.placeholder_format.type in BODY_TYPES for ph in layout.placeholders)


def resolve_layout_name(variant, slide, prs):
    names = [layout.name for layout in prs.slide_layouts]
    preferred = variant.get("rolePreferences", {}).get(slide.get("role"))
    if preferred and preferred in names:
        return preferred
    needs_body = bool(slide.get("bullets") or slide.get("body"))
    if needs_body:
        for layout in prs.slide_layouts:
            if _layout_has_body(layout):
                return layout.name
    return prs.slide_layouts[0].name


def _find_layout(prs, name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    return prs.slide_layouts[0]


def _set_text(placeholder, text=None, bullets=None):
    text_frame = placeholder.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    if bullets:
        text_frame.text = bullets[0]
        for bullet in bullets[1:]:
            paragraph = text_frame.add_paragraph()
            paragraph.text = bullet
    else:
        text_frame.text = text or ""


def _find_placeholder(slide, types):
    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.type in types:
            return placeholder
    return None


def fill_slide(slide, spec):
    label = spec.get("title") or spec.get("role") or "untitled slide"

    title = spec.get("title")
    if title and slide.shapes.title is not None:
        _set_text(slide.shapes.title, text=title)

    subtitle = spec.get("subtitle")
    if subtitle:
        placeholder = _find_placeholder(slide, (PP_PLACEHOLDER.SUBTITLE,))
        if placeholder is not None:
            _set_text(placeholder, text=subtitle)

    bullets = spec.get("bullets")
    body = spec.get("body")
    if bullets or body:
        placeholder = _find_placeholder(slide, BODY_TYPES)
        if placeholder is not None:
            if bullets is not None:
                if not isinstance(bullets, list):
                    raise ValueError(f'Slide "{label}" has a "bullets" field that isn\'t a list.')
                _set_text(placeholder, bullets=bullets)
            else:
                _set_text(placeholder, text=body)

    # Speaker notes — the talk track behind the slide (presenter view / printouts).
    notes = spec.get("notes")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _remove_all_slides(prs):
    slide_id_list = prs.slides._sldIdLst
    for slide_id in list(slide_id_list):
        slide_id_list.remove(slide_id)


def build_deck(config_path, variant_id, slides, output_path):
    if not isinstance(slides, list):
        raise ValueError(f"Expected a JSON array of slides, got {type(slides).__name__}.")
    config = load_config(config_path)
    variant = get_variant(config, variant_id)
    config_dir = os.path.dirname(os.path.abspath(config_path))
    template_path = os.path.join(config_dir, variant["template"])
    prs = Presentation(template_path)
    _remove_all_slides(prs)
    for spec in slides:
        layout_name = resolve_layout_name(variant, spec, prs)
        slide = prs.slides.add_slide(_find_layout(prs, layout_name))
        fill_slide(slide, spec)
    prs.save(output_path)
    return output_path


def _main(argv=None):
    import sys

    args = argv if argv is not None else sys.argv[1:]
    if len(args) != 4:
        print("Usage: python build_deck.py <deck-config.json> <slides.json> "
              "<variantId> <output.pptx>", file=sys.stderr)
        return 1
    config_path, slides_path, variant_id, output_path = args
    with open(slides_path) as f:
        slides = json.load(f)
    build_deck(config_path, variant_id, slides, output_path)
    print(f"Wrote {output_path}")
    return 0


if __name__ == "__main__":
    import sys

    try:
        sys.exit(_main())
    except Exception as exc:  # noqa: BLE001 - surface a clean message to the CLI user
        print(str(exc), file=sys.stderr)
        sys.exit(1)
