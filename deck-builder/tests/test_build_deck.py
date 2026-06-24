import os
import json
import shutil
import tempfile
import unittest

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from scripts.build_deck import (
    load_config,
    get_variant,
    resolve_layout_name,
    fill_slide,
    build_deck,
)


class TestConfigAndResolution(unittest.TestCase):
    def setUp(self):
        self.dir = tempfile.mkdtemp()
        self.addCleanup(shutil.rmtree, self.dir)
        self.template = os.path.join(self.dir, "templates", "v1.pptx")
        os.makedirs(os.path.dirname(self.template))
        Presentation().save(self.template)
        self.config_path = os.path.join(self.dir, "deck-config.json")
        with open(self.config_path, "w") as f:
            json.dump({"variants": [{
                "id": "v1", "name": "V1", "template": "templates/v1.pptx",
                "layouts": ["Title Slide", "Title and Content"],
                "rolePreferences": {"title": "Title Slide", "problem": "Title and Content"},
            }]}, f)

    def test_load_and_get_variant(self):
        cfg = load_config(self.config_path)
        v = get_variant(cfg, "v1")
        self.assertEqual(v["name"], "V1")

    def test_get_variant_missing_raises(self):
        cfg = load_config(self.config_path)
        with self.assertRaisesRegex(ValueError, 'No variant "nope"'):
            get_variant(cfg, "nope")

    def test_resolve_uses_role_preference(self):
        cfg = load_config(self.config_path)
        v = get_variant(cfg, "v1")
        prs = Presentation(self.template)
        name = resolve_layout_name(v, {"role": "title", "title": "Hi"}, prs)
        self.assertEqual(name, "Title Slide")

    def test_resolve_falls_back_to_body_layout_for_bullets(self):
        cfg = load_config(self.config_path)
        v = get_variant(cfg, "v1")
        v["rolePreferences"] = {}  # force fallback
        prs = Presentation(self.template)
        name = resolve_layout_name(v, {"role": "content", "bullets": ["a"]}, prs)
        layout = next(l for l in prs.slide_layouts if l.name == name)
        types = {ph.placeholder_format.type for ph in layout.placeholders}
        self.assertTrue(PP_PLACEHOLDER.BODY in types or PP_PLACEHOLDER.OBJECT in types)


class TestFillSlide(unittest.TestCase):
    def _slide(self, layout_name):
        self.prs = Presentation()
        layout = next(l for l in self.prs.slide_layouts if l.name == layout_name)
        return self.prs.slides.add_slide(layout)

    def test_fills_title_and_subtitle(self):
        slide = self._slide("Title Slide")
        fill_slide(slide, {"role": "title", "title": "Acme Review", "subtitle": "June 2026"})
        self.assertEqual(slide.shapes.title.text, "Acme Review")
        subs = [p for p in slide.placeholders
                if p.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE]
        self.assertEqual(subs[0].text, "June 2026")

    def test_fills_bullets_into_body(self):
        slide = self._slide("Title and Content")
        fill_slide(slide, {"role": "problem", "title": "Challenges",
                           "bullets": ["Alpha", "Beta", "Gamma"]})
        body = next(p for p in slide.placeholders
                    if p.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT))
        paragraph_texts = [para.text for para in body.text_frame.paragraphs]
        self.assertEqual(paragraph_texts, ["Alpha", "Beta", "Gamma"])

    def test_bullets_must_be_a_list(self):
        slide = self._slide("Title and Content")
        with self.assertRaisesRegex(ValueError, "isn't a list"):
            fill_slide(slide, {"role": "problem", "title": "X", "bullets": "not a list"})


class TestBuildDeck(unittest.TestCase):
    def setUp(self):
        self.dir = tempfile.mkdtemp()
        self.addCleanup(shutil.rmtree, self.dir)
        os.makedirs(os.path.join(self.dir, "templates"))
        Presentation().save(os.path.join(self.dir, "templates", "v1.pptx"))
        self.config_path = os.path.join(self.dir, "deck-config.json")
        with open(self.config_path, "w") as f:
            json.dump({"variants": [{
                "id": "v1", "name": "V1", "template": "templates/v1.pptx",
                "layouts": ["Title Slide", "Title and Content"],
                "rolePreferences": {"title": "Title Slide", "problem": "Title and Content"},
            }]}, f)

    def test_build_produces_pptx_with_right_layouts_and_text(self):
        out = os.path.join(self.dir, "out.pptx")
        slides = [
            {"role": "title", "title": "Acme Review", "subtitle": "June 2026"},
            {"role": "problem", "title": "Current Challenges", "bullets": ["Slow", "Manual"]},
        ]
        build_deck(self.config_path, "v1", slides, out)
        self.assertTrue(os.path.exists(out))
        prs = Presentation(out)
        self.assertEqual(len(prs.slides), 2)
        self.assertEqual(prs.slides[0].slide_layout.name, "Title Slide")
        self.assertEqual(prs.slides[0].shapes.title.text, "Acme Review")
        self.assertEqual(prs.slides[1].slide_layout.name, "Title and Content")

    def test_build_starts_from_an_empty_deck(self):
        tmpl = os.path.join(self.dir, "templates", "v1.pptx")
        prs = Presentation(tmpl)
        prs.slides.add_slide(prs.slide_layouts[0])  # pre-existing junk slide
        prs.save(tmpl)
        out = os.path.join(self.dir, "out2.pptx")
        build_deck(self.config_path, "v1", [{"role": "title", "title": "Only Me"}], out)
        self.assertEqual(len(Presentation(out).slides), 1)


if __name__ == "__main__":
    unittest.main()
