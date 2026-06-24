import os
import json
import shutil
import tempfile
import unittest
import importlib.util

from pptx import Presentation

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))


def _load(module_name, file_path):
    spec = importlib.util.spec_from_file_location(module_name, file_path)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


# Both skills expose a `scripts` package; load each by explicit file path under a unique
# module name so the two never collide.
setup_mod = _load("dts_extract_layouts",
                  os.path.join(ROOT, "deck-theme-setup", "scripts", "extract_layouts.py"))
build_mod = _load("db_build_deck",
                  os.path.join(ROOT, "deck-builder", "scripts", "build_deck.py"))


class TestSeam(unittest.TestCase):
    """Real deck-theme-setup output must feed deck-builder unchanged."""

    def _example_deck(self, dir_):
        prs = Presentation()
        layouts = {l.name: l for l in prs.slide_layouts}
        prs.slides.add_slide(layouts["Title Slide"]).shapes.title.text = "Welcome"
        prs.slides.add_slide(layouts["Title and Content"]).shapes.title.text = "Current Challenges"
        path = os.path.join(dir_, "example.pptx")
        prs.save(path)
        return path

    def test_setup_learn_then_build(self):
        d = tempfile.mkdtemp()
        self.addCleanup(shutil.rmtree, d)
        template = os.path.join(d, "src.pptx")
        Presentation().save(template)
        config = os.path.join(d, "deck-config.json")

        setup_mod.run_setup(template, "v1", "V1", config, os.path.join(d, "templates"))
        setup_mod.run_learn(self._example_deck(d), "v1", config)

        out = os.path.join(d, "out.pptx")
        slides = [
            {"role": "title", "title": "Acme", "subtitle": "2026"},
            {"role": "problem", "title": "Current Challenges", "bullets": ["A", "B"]},
        ]
        build_mod.build_deck(config, "v1", slides, out)

        prs = Presentation(out)
        self.assertEqual(len(prs.slides), 2)
        self.assertEqual(prs.slides[0].slide_layout.name, "Title Slide")
        self.assertEqual(prs.slides[1].slide_layout.name, "Title and Content")
        self.assertEqual(prs.slides[0].shapes.title.text, "Acme")


if __name__ == "__main__":
    unittest.main()
